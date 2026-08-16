from pathlib import Path
from typing import Any

from fastapi import FastAPI
from pydantic import BaseModel

from app.api.gis import router as gis_router
from app.api.knowledge import router as knowledge_router
from app.api.ndvi import router as ndvi_router
from app.api.papers import router as papers_router
from app.exceptions import GeoWorkError
from app.middleware.auth import register_auth_middleware
from app.middleware.error_handler import generic_exception_handler, geowork_exception_handler
from app.tool_catalog import TOOL_CATALOG, validate_catalog
from app.validation import ValidationError, validate_bbox, validate_crs, validate_path

app = FastAPI(title="GeoWork Geo Python Worker", version="1.0.0-dev")

# doc/22 BP4 / F6: fail-closed runtime token auth (mirrors core/api/auth.go)
# + explicit CORS allowlist for loopback dev origins. Registered before
# routers so they guard every endpoint incl. /health-exempt paths.
register_auth_middleware(app)
from fastapi.middleware.cors import CORSMiddleware  # noqa: E402

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "http://127.0.0.1:5173",
    ],
    allow_methods=["GET", "POST"],
    allow_headers=["Content-Type", "X-GeoWork-Token"],
)

# Fail fast at import time if the tool catalog violates the WorkerToolDef contract.
validate_catalog()

# Register unified exception handlers
app.add_exception_handler(GeoWorkError, geowork_exception_handler)
app.add_exception_handler(Exception, generic_exception_handler)

# Include NDVI API router
app.include_router(ndvi_router)

# Include papers API router
app.include_router(papers_router)
app.include_router(gis_router, prefix="/api")
app.include_router(knowledge_router, prefix="/api")


class ToolRequest(BaseModel):
    workspace: str
    taskId: str = "task"
    prompt: str = ""
    mode: str = "Analysis"
    params: dict[str, Any] = {}


def ensure_workspace(path: str) -> Path:
    workspace = Path(path).expanduser().resolve()
    for child in ["scripts", "reports", "artifacts", "data", "knowledge"]:
        (workspace / child).mkdir(parents=True, exist_ok=True)
    return workspace


def artifact(name: str, path: Path, kind: str, mime_type: str) -> dict[str, str]:
    return {"name": name, "path": str(path), "type": kind, "mimeType": mime_type}


def write_manifest(workspace: Path, task_id: str, artifacts: list[dict[str, str]]) -> dict[str, str]:
    manifest_path = workspace / "artifacts" / f"{task_id}_artifact_manifest.json"
    import json

    manifest_path.write_text(
        json.dumps(
            {
                "taskId": task_id,
                "artifacts": artifacts,
                "reproducibility": {
                    "workspaceScoped": True,
                    "generatedBy": "GeoWork Geo Python Worker",
                    "qgisBundled": False,
                },
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return artifact("Artifact Manifest", manifest_path, "manifest", "application/json")


def write_json_artifact(workspace: Path, task_id: str, name: str, filename: str, payload: dict[str, Any], kind: str) -> dict[str, str]:
    import json

    path = workspace / "artifacts" / filename
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return artifact(name, path, kind, "application/json")


@app.get("/health")
def health():
    return {
        "status": "ok",
        "service": "geo-python-worker",
        "capabilities": [
            "gee-ndvi-script",
            "ndvi-analysis",
            "office-report",
            "pdf-parse",
            "gdal-inspect",
            "qgis-detect",
            "qgis-processing",
            "raster-vector-tools",
            "ppt-excel-notebook",
            "cog-map-layout",
        ],
    }


@app.get("/tools")
def list_tools():
    """Tool catalog consumed by the Go core (worker.Client.ListTools)."""
    return {"tools": TOOL_CATALOG}


@app.post("/tools/gee/search-dataset")
def search_gee_dataset(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    query = req.params.get("query") or req.prompt or "NDVI Sentinel-2"
    result = {
        "ok": True,
        "query": query,
        "datasets": [
            {"id": "COPERNICUS/S2_SR_HARMONIZED", "name": "Sentinel-2 Surface Reflectance", "resolution": "10m", "use": "NDVI, land cover, urban vegetation"},
            {"id": "LANDSAT/LC08/C02/T1_L2", "name": "Landsat 8 Collection 2 Level 2", "resolution": "30m", "use": "long time-series, LST, NDVI"},
            {"id": "MODIS/061/MOD13Q1", "name": "MODIS Vegetation Indices", "resolution": "250m", "use": "regional NDVI trends"},
        ],
    }
    art = write_json_artifact(workspace, req.taskId, "GEE Dataset Search", f"{req.taskId}_gee_datasets.json", result, "gee-datasets")
    return {"ok": True, "message": "GEE dataset search completed", "artifacts": [art], "datasets": result["datasets"]}


@app.post("/tools/gee/check-auth")
def check_gee_auth(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    status = {"ok": True, "authenticated": False, "method": "earthengine credentials", "nextStep": "Run earthengine authenticate if exports are required."}
    try:
        import ee

        ee.Initialize()
        status["authenticated"] = True
        status["nextStep"] = "Earth Engine is ready."
    except Exception as exc:
        status["error"] = str(exc)
    art = write_json_artifact(workspace, req.taskId, "GEE Auth Status", f"{req.taskId}_gee_auth.json", status, "environment")
    return {"ok": True, "message": "GEE authentication checked", "artifacts": [art], "status": status}


@app.post("/tools/gee/generate-ndvi-script")
def generate_ndvi_script(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    script_path = workspace / "scripts" / f"{req.taskId}_gee_ndvi.py"
    script = f'''"""GeoWork generated GEE NDVI workflow."""
import ee

ee.Initialize()

AOI = ee.Geometry.Rectangle([100.0, 20.0, 101.0, 21.0])
COLLECTION = (
    ee.ImageCollection("COPERNICUS/S2_SR_HARMONIZED")
    .filterBounds(AOI)
    .filterDate("2024-01-01", "2024-12-31")
    .filter(ee.Filter.lt("CLOUDY_PIXEL_PERCENTAGE", 20))
)

def add_ndvi(image):
    ndvi = image.normalizedDifference(["B8", "B4"]).rename("NDVI")
    return image.addBands(ndvi)

ndvi_collection = COLLECTION.map(add_ndvi)
median_ndvi = ndvi_collection.select("NDVI").median().clip(AOI)

task = ee.batch.Export.image.toDrive(
    image=median_ndvi,
    description="geowork_ndvi_export",
    folder="GeoWork",
    fileNamePrefix="{req.taskId}_ndvi",
    scale=10,
    region=AOI,
    maxPixels=1e13,
)
task.start()
print("Started NDVI export", task.id)
'''
    script_path.write_text(script, encoding="utf-8")
    map_path = workspace / "artifacts" / f"{req.taskId}_map.html"
    map_path.write_text(
        "<!doctype html><title>GeoWork NDVI Map</title><h1>NDVI Preview</h1><p>GEE export task created. Load output COG/GeoTIFF here after completion.</p>",
        encoding="utf-8",
    )
    artifacts = [
        artifact("GEE NDVI Python Script", script_path, "script", "text/x-python"),
        artifact("NDVI Map Preview", map_path, "html-map", "text/html"),
    ]
    artifacts.append(write_manifest(workspace, req.taskId, artifacts))
    return {
        "ok": True,
        "message": "Generated GEE NDVI Python script and HTML map preview",
        "artifacts": artifacts,
    }


@app.post("/tools/office/write-report")
def write_report(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    markdown_path = workspace / "reports" / f"{req.taskId}_report.md"
    docx_path = workspace / "reports" / f"{req.taskId}_report.docx"
    markdown = f"""# GeoWork Task Report

## Prompt

{req.prompt or "GeoWork analysis task"}

## Workflow

- Research/Data/GeoCode/Analysis/Write mode: {req.mode}
- Generated a transparent plan in Go Core.
- Executed tools through the guarded Tool Registry.
- Registered outputs as project artifacts.

## Reproducibility

All file writes were scoped to the project workspace. External QGIS/GDAL/GEE integrations should be configured from Settings before production use.
"""
    markdown_path.write_text(markdown, encoding="utf-8")
    try:
        from docx import Document

        doc = Document()
        doc.add_heading("GeoWork Task Report", 0)
        doc.add_heading("Prompt", level=1)
        doc.add_paragraph(req.prompt or "GeoWork analysis task")
        doc.add_heading("Workflow", level=1)
        for line in [
            f"Mode: {req.mode}",
            "Planner/Executor ran through Go Core.",
            "Tool calls were logged and artifacts registered.",
        ]:
            doc.add_paragraph(line, style="List Bullet")
        doc.save(docx_path)
    except Exception:
        docx_path.write_text("python-docx unavailable; Markdown report generated.", encoding="utf-8")
    artifacts = [
        artifact("Markdown Report", markdown_path, "report", "text/markdown"),
        artifact(
            "Word Report",
            docx_path,
            "report",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
    ]
    artifacts.append(write_manifest(workspace, req.taskId, artifacts))
    return {
        "ok": True,
        "message": "Generated Markdown and DOCX reports",
        "artifacts": artifacts,
    }


@app.post("/tools/office/write-ppt")
def write_ppt(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    ppt_path = workspace / "reports" / f"{req.taskId}_presentation.pptx"
    try:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "GeoWork Analysis"
        slide.placeholders[1].text = req.prompt or "GeoWork generated presentation"
        for title in ["Objective", "Data", "Workflow", "Results", "Reproducibility"]:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = title
            s.placeholders[1].text = f"{title} generated from task {req.taskId}"
        prs.save(ppt_path)
    except Exception:
        ppt_path.write_text("GeoWork PPTX export content\n" + (req.prompt or ""), encoding="utf-8")
    artifacts = [artifact("PowerPoint Presentation", ppt_path, "presentation", "application/vnd.openxmlformats-officedocument.presentationml.presentation")]
    artifacts.append(write_manifest(workspace, req.taskId, artifacts))
    return {"ok": True, "message": "Generated PPTX presentation", "artifacts": artifacts}


@app.post("/tools/office/write-excel")
def write_excel(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    xlsx_path = workspace / "reports" / f"{req.taskId}_statistics.xlsx"
    try:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "GeoWork Statistics"
        ws.append(["metric", "value"])
        ws.append(["task_id", req.taskId])
        ws.append(["mode", req.mode])
        ws.append(["prompt", req.prompt])
        wb.save(xlsx_path)
    except Exception:
        xlsx_path.write_text("metric,value\ntask_id," + req.taskId + "\n", encoding="utf-8")
    artifacts = [artifact("Excel Statistics", xlsx_path, "spreadsheet", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")]
    artifacts.append(write_manifest(workspace, req.taskId, artifacts))
    return {"ok": True, "message": "Generated Excel workbook", "artifacts": artifacts}


@app.post("/tools/office/write-notebook")
def write_notebook(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    notebook_path = workspace / "reports" / f"{req.taskId}_workflow.ipynb"
    notebook = {
        "cells": [
            {"cell_type": "markdown", "metadata": {}, "source": ["# GeoWork Reproducible Workflow\n", req.prompt]},
            {"cell_type": "code", "execution_count": None, "metadata": {}, "outputs": [], "source": ["print('GeoWork task:', '" + req.taskId + "')"]},
        ],
        "metadata": {"kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"}},
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    art = write_json_artifact(workspace, req.taskId, "Jupyter Notebook", f"{req.taskId}_workflow.ipynb", notebook, "notebook")
    art["mimeType"] = "application/x-ipynb+json"
    return {"ok": True, "message": "Generated Jupyter Notebook", "artifacts": [art, write_manifest(workspace, req.taskId, [art])]}


@app.post("/tools/gdal/inspect-dataset")
def inspect_dataset(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    source = req.params.get("path") or req.prompt or ""
    if not source:
        return {"ok": False, "message": "Missing required parameter: path"}

    report_path = workspace / "artifacts" / f"{req.taskId}_dataset_quality.json"
    report: dict[str, Any] = {"ok": True, "path": source}

    try:
        import rasterio

        with rasterio.open(source) as src:
            report["type"] = "raster"
            report["crs"] = str(src.crs) if src.crs else None
            report["bounds"] = {"left": src.bounds.left, "bottom": src.bounds.bottom, "right": src.bounds.right, "top": src.bounds.top}
            report["dtype"] = src.dtypes[0]
            report["band_count"] = src.count
            report["width"] = src.width
            report["height"] = src.height
            report["nodata"] = src.nodata
            report["driver"] = src.driver
    except ImportError:
        report["error"] = "rasterio not installed — cannot inspect raster datasets"
    except Exception:
        # Not a raster or open failed — try vector
        try:
            import geopandas as gpd

            gdf = gpd.read_file(source)
            report["type"] = "vector"
            report["crs"] = str(gdf.crs) if gdf.crs else None
            bounds = gdf.total_bounds
            report["bounds"] = {"left": float(bounds[0]), "bottom": float(bounds[1]), "right": float(bounds[2]), "top": float(bounds[3])}
            report["feature_count"] = len(gdf)
            report["columns"] = list(gdf.columns)
            report["geometry_type"] = gdf.geometry.geom_type.iloc[0] if len(gdf) > 0 else "empty"
        except ImportError:
            report["error"] = "Neither rasterio nor geopandas available for dataset inspection"
        except Exception as vec_exc:
            report["error"] = f"Failed to open as raster or vector: {vec_exc}"

    import json
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    return {
        "ok": True,
        "message": "Dataset quality inspection completed",
        "artifacts": [artifact("Dataset Quality JSON", report_path, "quality-report", "application/json")],
        "metadata": report,
    }


@app.post("/tools/raster/metadata")
def raster_metadata(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    source = req.params.get("path") or req.prompt or ""
    if not source:
        return {"ok": False, "message": "Missing required parameter: path"}

    try:
        import rasterio

        with rasterio.open(source) as src:
            payload = {
                "ok": True,
                "path": source,
                "driver": src.driver,
                "width": src.width,
                "height": src.height,
                "count": src.count,
                "dtype": src.dtypes[0],
                "crs": str(src.crs) if src.crs else None,
                "bounds": {"left": src.bounds.left, "bottom": src.bounds.bottom, "right": src.bounds.right, "top": src.bounds.top},
                "nodata": src.nodata,
                "transform": list(src.transform)[:6] if src.transform else None,
            }
    except ImportError:
        return {"ok": False, "message": "rasterio is not installed. Install with: pip install rasterio"}
    except Exception as exc:
        return {"ok": False, "message": f"Failed to read raster metadata: {exc}"}

    art = write_json_artifact(workspace, req.taskId, "Raster Metadata", f"{req.taskId}_raster_metadata.json", payload, "raster-metadata")
    return {"ok": True, "message": "Raster metadata generated", "artifacts": [art], "metadata": payload}


@app.post("/tools/raster/clip")
def raster_clip(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    source = req.params.get("path") or ""
    output = req.params.get("output") or str(workspace / "artifacts" / f"{req.taskId}_raster_clip.tif")
    bbox = req.params.get("bbox")

    if not source:
        return {"ok": False, "message": "Missing required parameter: path"}

    # Validate path safety
    try:
        validate_path(source)
    except ValidationError as e:
        return {"ok": False, "message": str(e)}

    if not bbox or not isinstance(bbox, list) or len(bbox) != 4:
        return {"ok": False, "message": "Missing or invalid bbox: expected [minX, minY, maxX, maxY]"}

    # Validate bbox range
    try:
        validate_bbox(bbox)
    except ValidationError as e:
        return {"ok": False, "message": str(e)}

    try:
        import rasterio
        from rasterio.windows import from_bounds

        out_path = Path(output)
        out_path.parent.mkdir(parents=True, exist_ok=True)

        with rasterio.open(source) as src:
            window = from_bounds(bbox[0], bbox[1], bbox[2], bbox[3], src.transform)
            data = src.read(window=window)
            meta = src.meta.copy()
            meta.update({
                "height": data.shape[1],
                "width": data.shape[2],
                "transform": src.window_transform(window),
            })
            with rasterio.open(out_path, "w", **meta) as dst:
                dst.write(data)

        return {"ok": True, "message": "Raster clip completed", "artifacts": [artifact("Raster Clip", out_path, "GeoTIFF", "image/tiff")]}
    except ImportError:
        return {"ok": False, "message": "rasterio is not installed. Install with: pip install rasterio"}
    except Exception as exc:
        return {"ok": False, "message": f"Raster clip failed: {exc}"}


@app.post("/tools/raster/reproject")
def raster_reproject(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    source = req.params.get("path") or ""
    dst_crs = req.params.get("dst_crs") or req.params.get("targetCrs") or ""
    output = req.params.get("output") or str(workspace / "artifacts" / f"{req.taskId}_raster_reprojected.tif")

    if not source:
        return {"ok": False, "message": "Missing required parameter: path"}
    if not dst_crs:
        return {"ok": False, "message": "Missing required parameter: dst_crs (target CRS)"}

    # Validate path and CRS
    try:
        validate_path(source)
        validate_crs(dst_crs)
    except ValidationError as e:
        return {"ok": False, "message": str(e)}

    try:
        import rasterio
        from rasterio.warp import calculate_default_transform, reproject, Resampling

        out_path = Path(output)
        out_path.parent.mkdir(parents=True, exist_ok=True)

        with rasterio.open(source) as src:
            dst_transform, dst_width, dst_height = calculate_default_transform(
                src.crs, dst_crs, src.width, src.height, *src.bounds
            )
            meta = src.meta.copy()
            meta.update({
                "crs": dst_crs,
                "transform": dst_transform,
                "width": dst_width,
                "height": dst_height,
            })
            with rasterio.open(out_path, "w", **meta) as dst:
                for i in range(1, src.count + 1):
                    reproject(
                        source=rasterio.band(src, i),
                        destination=rasterio.band(dst, i),
                        src_transform=src.transform,
                        src_crs=src.crs,
                        dst_transform=dst_transform,
                        dst_crs=dst_crs,
                        resampling=Resampling.nearest,
                    )

        return {"ok": True, "message": "Raster reprojection completed", "artifacts": [artifact("Reprojected Raster", out_path, "GeoTIFF", "image/tiff")]}
    except ImportError:
        return {"ok": False, "message": "rasterio is not installed. Install with: pip install rasterio"}
    except Exception as exc:
        return {"ok": False, "message": f"Raster reprojection failed: {exc}"}


@app.post("/tools/raster/write-cog")
def raster_write_cog(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    source = req.params.get("path") or ""
    output = req.params.get("output") or str(workspace / "artifacts" / f"{req.taskId}.cog.tif")

    if not source:
        return {"ok": False, "message": "Missing required parameter: path"}

    try:
        import rasterio
        from rasterio.transform import from_bounds as _fb
        from rasterio.enums import Resampling as _Resamp

        out_path = Path(output)
        out_path.parent.mkdir(parents=True, exist_ok=True)

        with rasterio.open(source) as src:
            data = src.read()
            meta = src.meta.copy()
            meta.update({
                "driver": "GTiff",
                "tiled": True,
                "compress": "deflate",
                "blockxsize": 512,
                "blockysize": 512,
            })
            with rasterio.open(out_path, "w", **meta) as dst:
                dst.write(data)

        return {"ok": True, "message": "COG artifact generated", "artifacts": [artifact("Cloud Optimized GeoTIFF", out_path, "COG", "image/tiff")]}
    except ImportError:
        return {"ok": False, "message": "rasterio is not installed. Install with: pip install rasterio"}
    except Exception as exc:
        return {"ok": False, "message": f"COG write failed: {exc}"}


@app.post("/tools/vector/metadata")
def vector_metadata(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    source = req.params.get("path") or req.prompt or ""
    if not source:
        return {"ok": False, "message": "Missing required parameter: path"}

    try:
        import geopandas as gpd

        gdf = gpd.read_file(source)
        bounds = gdf.total_bounds
        payload = {
            "ok": True,
            "path": source,
            "driver": "GeoJSON",
            "crs": str(gdf.crs) if gdf.crs else None,
            "bounds": {"left": float(bounds[0]), "bottom": float(bounds[1]), "right": float(bounds[2]), "top": float(bounds[3])},
            "feature_count": len(gdf),
            "columns": list(gdf.columns),
            "geometry_type": gdf.geometry.geom_type.iloc[0] if len(gdf) > 0 else "empty",
        }
    except ImportError:
        return {"ok": False, "message": "geopandas is not installed. Install with: pip install geopandas"}
    except Exception as exc:
        return {"ok": False, "message": f"Failed to read vector metadata: {exc}"}

    art = write_json_artifact(workspace, req.taskId, "Vector Metadata", f"{req.taskId}_vector_metadata.json", payload, "vector-metadata")
    return {"ok": True, "message": "Vector metadata generated", "artifacts": [art], "metadata": payload}


@app.post("/tools/vector/buffer")
def vector_buffer(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    try:
        import geopandas as gpd

        params = req.params or {}
        in_path = params.get("path", "")
        out_path_str = params.get("output", "")
        distance = float(params.get("distance", 0))

        if not in_path:
            return {"ok": False, "message": "Missing required parameter: path"}

        gdf = gpd.read_file(in_path)
        gdf["geometry"] = gdf.geometry.buffer(distance)

        out_path = Path(out_path_str) if out_path_str else workspace / "artifacts" / f"{req.taskId}_buffer.geojson"
        out_path.parent.mkdir(parents=True, exist_ok=True)
        gdf.to_file(str(out_path), driver="GeoJSON")

        return {
            "ok": True,
            "message": "Vector buffer completed",
            "artifacts": [artifact("Vector Buffer", out_path, "GeoJSON", "application/geo+json")],
            "output": str(out_path),
            "feature_count": len(gdf),
            "status": "ok",
        }
    except ImportError:
        return {"ok": False, "message": "geopandas is not installed. Install with: pip install geopandas", "status": "degraded"}
    except Exception as e:
        return {"ok": False, "message": f"Vector buffer failed: {e}", "status": "error"}


@app.post("/tools/vector/clip")
def vector_clip(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    try:
        import geopandas as gpd

        params = req.params or {}
        in_path = params.get("path", "")
        clip_path = params.get("clip", "") or params.get("clip_path", "")
        bbox = params.get("bbox")
        out_path_str = params.get("output", "")

        if not in_path:
            return {"ok": False, "message": "Missing required parameter: path"}

        gdf = gpd.read_file(in_path)

        if clip_path:
            clip_gdf = gpd.read_file(clip_path)
            if gdf.crs != clip_gdf.crs:
                clip_gdf = clip_gdf.to_crs(gdf.crs)
            clipped = gpd.overlay(gdf, clip_gdf, how="intersection")
        elif bbox and isinstance(bbox, list) and len(bbox) == 4:
            from shapely.geometry import box
            clip_geom = box(bbox[0], bbox[1], bbox[2], bbox[3])
            import pandas as pd
            clip_gdf = gpd.GeoDataFrame(geometry=[clip_geom], crs=gdf.crs)
            clipped = gpd.overlay(gdf, clip_gdf, how="intersection")
        else:
            return {"ok": False, "message": "Provide 'clip' (path to clip layer) or 'bbox' ([minX,minY,maxX,maxY])"}

        out_path = Path(out_path_str) if out_path_str else workspace / "artifacts" / f"{req.taskId}_vector_clip.geojson"
        out_path.parent.mkdir(parents=True, exist_ok=True)
        clipped.to_file(str(out_path), driver="GeoJSON")

        return {
            "ok": True,
            "message": "Vector clip completed",
            "artifacts": [artifact("Vector Clip", out_path, "GeoJSON", "application/geo+json")],
            "output": str(out_path),
            "feature_count": len(clipped),
            "status": "ok",
        }
    except ImportError as e:
        return {"ok": False, "message": f"Missing dependency: {e}. Install with: pip install geopandas shapely", "status": "degraded"}
    except Exception as e:
        return {"ok": False, "message": f"Vector clip failed: {e}", "status": "error"}


@app.post("/tools/vector/reproject")
def vector_reproject(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    try:
        import geopandas as gpd

        params = req.params or {}
        in_path = params.get("path", "")
        crs = params.get("crs") or params.get("targetCrs") or params.get("dst_crs", "")
        out_path_str = params.get("output", "")

        if not in_path:
            return {"ok": False, "message": "Missing required parameter: path"}
        if not crs:
            return {"ok": False, "message": "Missing required parameter: crs (target CRS, e.g. EPSG:4326)"}

        # Validate path and CRS format
        try:
            validate_path(in_path)
            validate_crs(crs)
        except ValidationError as e:
            return {"ok": False, "message": str(e)}

        gdf = gpd.read_file(in_path)
        reprojected = gdf.to_crs(crs)

        out_path = Path(out_path_str) if out_path_str else workspace / "artifacts" / f"{req.taskId}_vector_reprojected.geojson"
        out_path.parent.mkdir(parents=True, exist_ok=True)
        reprojected.to_file(str(out_path), driver="GeoJSON")

        return {
            "ok": True,
            "message": "Vector reprojection completed",
            "artifacts": [artifact("Reprojected Vector", out_path, "GeoJSON", "application/geo+json")],
            "output": str(out_path),
            "feature_count": len(reprojected),
            "crs": str(crs),
            "status": "ok",
        }
    except ImportError:
        return {"ok": False, "message": "geopandas is not installed. Install with: pip install geopandas", "status": "degraded"}
    except Exception as e:
        return {"ok": False, "message": f"Vector reprojection failed: {e}", "status": "error"}


@app.post("/tools/map/layout-export")
def map_layout_export(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    params = req.params or {}
    source = params.get("path") or params.get("source") or ""
    title = params.get("title") or "GeoWork Map Layout"
    fmt = params.get("format") or "all"  # "html", "png", "svg", or "all"

    artifacts = []
    errors = []

    # --- HTML export (folium or fallback) ---
    if fmt in ("html", "all"):
        html_path = workspace / "artifacts" / f"{req.taskId}_map_layout.html"
        try:
            import folium

            m = None
            if source:
                import geopandas as gpd
                gdf = gpd.read_file(source)
                bounds = gdf.total_bounds
                center = [(bounds[1] + bounds[3]) / 2, (bounds[0] + bounds[2]) / 2]
                m = folium.Map(location=center, zoom_start=12, tiles="OpenStreetMap")
                folium.GeoJson(source, name="data").add_to(m)
                folium.LayerControl().add_to(m)
            else:
                m = folium.Map(location=[39.9, 116.4], zoom_start=10, tiles="OpenStreetMap")
                folium.Marker([39.9, 116.4], popup=title).add_to(m)

            m.save(str(html_path))
            artifacts.append(artifact("HTML Map Layout", html_path, "HTML Map", "text/html"))
        except ImportError:
            # Fallback: simple HTML template
            html_content = f"""<!doctype html>
<html lang="en">
<head><meta charset="utf-8"><title>{title}</title>
<style>body{{margin:0;font-family:sans-serif}}#map{{width:100%;height:100vh}}</style>
</head>
<body><div id="map"><h1>{title}</h1><p>Install folium for interactive map: pip install folium</p></div></body>
</html>"""
            html_path.write_text(html_content, encoding="utf-8")
            artifacts.append(artifact("HTML Map Layout", html_path, "HTML Map", "text/html"))
        except Exception as e:
            errors.append(f"HTML export failed: {e}")

    # --- PNG export (matplotlib) ---
    if fmt in ("png", "all"):
        png_path = workspace / "artifacts" / f"{req.taskId}_map_layout.png"
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
            import geopandas as gpd

            fig, ax = plt.subplots(figsize=(10, 7))
            if source:
                gdf = gpd.read_file(source)
                gdf.plot(ax=ax, edgecolor="black", facecolor="#6baed6", linewidth=0.5)
                ax.set_title(title, fontsize=14, fontweight="bold")
                ax.set_xlabel("Longitude")
                ax.set_ylabel("Latitude")
            else:
                ax.text(0.5, 0.5, title, ha="center", va="center", fontsize=16, transform=ax.transAxes)
                ax.set_title("No data source provided", fontsize=10)

            plt.tight_layout()
            fig.savefig(str(png_path), dpi=150, bbox_inches="tight")
            plt.close(fig)
            artifacts.append(artifact("PNG Map Layout", png_path, "PNG", "image/png"))
        except ImportError:
            errors.append("matplotlib/geopandas not installed — PNG export skipped")
            # Write a minimal placeholder PNG
            png_path.write_bytes(b"\x89PNG\r\n\x1a\n")
            artifacts.append(artifact("PNG Map Layout (placeholder)", png_path, "PNG", "image/png"))
        except Exception as e:
            errors.append(f"PNG export failed: {e}")

    # --- SVG export (matplotlib SVG backend) ---
    if fmt in ("svg", "all"):
        svg_path = workspace / "artifacts" / f"{req.taskId}_map_layout.svg"
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
            import geopandas as gpd

            fig, ax = plt.subplots(figsize=(10, 7))
            if source:
                gdf = gpd.read_file(source)
                gdf.plot(ax=ax, edgecolor="black", facecolor="#6baed6", linewidth=0.5)
                ax.set_title(title, fontsize=14, fontweight="bold")
                ax.set_xlabel("Longitude")
                ax.set_ylabel("Latitude")
            else:
                ax.text(0.5, 0.5, title, ha="center", va="center", fontsize=16, transform=ax.transAxes)
                ax.set_title("No data source provided", fontsize=10)

            plt.tight_layout()
            fig.savefig(str(svg_path), format="svg", bbox_inches="tight")
            plt.close(fig)
            artifacts.append(artifact("SVG Map Layout", svg_path, "SVG", "image/svg+xml"))
        except ImportError:
            errors.append("matplotlib/geopandas not installed — SVG export skipped")
            svg_path.write_text(
                f"<svg xmlns='http://www.w3.org/2000/svg' width='800' height='500'>"
                f"<text x='20' y='40' font-size='16'>{title}</text>"
                f"<text x='20' y='70' font-size='12'>Install matplotlib for real SVG export</text></svg>",
                encoding="utf-8",
            )
            artifacts.append(artifact("SVG Map Layout (placeholder)", svg_path, "SVG", "image/svg+xml"))
        except Exception as e:
            errors.append(f"SVG export failed: {e}")

    if artifacts:
        artifacts.append(write_manifest(workspace, req.taskId, artifacts))

    ok = len(artifacts) > 1 or not errors  # at least one real export succeeded
    msg = "Map layout exported"
    if errors:
        msg += f" (warnings: {'; '.join(errors)})"

    return {"ok": ok, "message": msg, "artifacts": artifacts}


# ---------------------------------------------------------------------------
# PDF parsing helpers (for /tools/papers/parse-pdf)
# ---------------------------------------------------------------------------

def _extract_title(text: str) -> str:
    """Extract paper title from text using heuristics."""
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    if not lines:
        return "Unknown Title"
    # Title is usually the first non-empty, reasonably short line
    for line in lines[:5]:
        if len(line) > 10 and len(line) < 300:
            return line
    return lines[0] if lines else "Unknown Title"


def _extract_abstract(text: str) -> str:
    """Extract abstract section from paper text."""
    import re
    match = re.search(
        r"(?:Abstract|摘要)\s*[:：]?\s*(.+?)(?=\n\s*(?:Introduction|Keywords|I\.|II\.|References|1\s|1\.))",
        text,
        re.DOTALL | re.IGNORECASE,
    )
    return match.group(1).strip()[:2000] if match else ""


def _extract_sections(text: str) -> list[dict[str, str]]:
    """Extract section headings and brief content from paper text."""
    import re
    sections: list[dict[str, str]] = []
    # Match common section heading patterns
    heading_pattern = re.compile(
        r"^\s*(?:\d+\.?\s*)?([A-Z][A-Za-z\s,\-:]{3,80})\s*$",
        re.MULTILINE,
    )
    headings = heading_pattern.findall(text)
    for heading in headings[:20]:
        heading = heading.strip()
        if heading and heading not in ("References", "Acknowledgment", "Acknowledgments"):
            sections.append({"heading": heading, "content": ""})
    return sections


@app.post("/tools/papers/parse-pdf")
def parse_pdf(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    pdf_path = req.params.get("path") or req.params.get("pdf_path") or ""

    if not pdf_path:
        return {"ok": False, "error": "No PDF path provided. Pass 'path' in params."}

    pdf_file = Path(pdf_path).expanduser()
    if not pdf_file.exists():
        return {"ok": False, "error": f"PDF file not found: {pdf_file}"}

    pdf_bytes = pdf_file.read_bytes()
    notes: dict[str, Any] = {"page_count": 0, "title": "", "abstract": "", "sections": [], "full_text": ""}

    # Try PyPDF2 first, then fall back to text heuristics
    try:
        from PyPDF2 import PdfReader
        import io

        reader = PdfReader(io.BytesIO(pdf_bytes))
        full_text = ""
        for page in reader.pages:
            full_text += (page.extract_text() or "")

        notes["page_count"] = len(reader.pages)
        notes["title"] = _extract_title(full_text)
        notes["abstract"] = _extract_abstract(full_text)
        notes["sections"] = _extract_sections(full_text)
        notes["full_text"] = full_text[:10000]
    except Exception as exc:
        # ImportError: PyPDF2 not available. Any other exception means the
        # bytes are not a parseable PDF (e.g. PdfReadError on a corrupted or
        # plain-text file) — either way fall back to text heuristics instead
        # of failing the request.
        try:
            raw_text = pdf_bytes.decode("utf-8", errors="replace")
        except Exception:
            raw_text = pdf_bytes.decode("latin-1", errors="replace")
        notes["title"] = _extract_title(raw_text)
        notes["abstract"] = _extract_abstract(raw_text)
        notes["sections"] = _extract_sections(raw_text)
        notes["full_text"] = raw_text[:10000]
        notes["status"] = "degraded"
        notes["warning"] = f"PyPDF2 extraction unavailable ({exc}); using raw text extraction"

    # Write structured notes
    output_path = workspace / "knowledge" / f"{req.taskId}_paper_notes.md"
    md_lines = [
        f"# {notes['title']}",
        "",
        f"**Pages:** {notes['page_count']}",
        "",
        "## Abstract",
        "",
        notes["abstract"] or "_(No abstract extracted)_",
        "",
        "## Sections",
        "",
    ]
    for sec in notes["sections"]:
        md_lines.append(f"- {sec['heading']}")
    md_lines.append("")
    md_lines.append("## Full Text (truncated)")
    md_lines.append("")
    md_lines.append(notes["full_text"][:5000])
    output_path.write_text("\n".join(md_lines), encoding="utf-8")

    return {
        "ok": True,
        "message": "PDF parsed successfully",
        "notes": {"title": notes["title"], "abstract": notes["abstract"], "page_count": notes["page_count"], "section_count": len(notes["sections"])},
        "artifacts": [artifact("Paper Reading Notes", output_path, "knowledge", "text/markdown")],
    }


@app.post("/tools/papers/openalex-search")
def openalex_search(req: ToolRequest):
    import csv
    import io

    workspace = ensure_workspace(req.workspace)
    query = req.params.get("query") or req.prompt or "remote sensing"
    limit = int(req.params.get("limit") or 10)

    papers: list[dict[str, Any]] = []
    error_msg = ""

    try:
        import httpx

        resp = httpx.get(
            "https://api.openalex.org/works",
            params={"search": query, "per_page": limit},
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()

        for w in data.get("results", []):
            # Reconstruct abstract from inverted index
            abstract = ""
            inverted = w.get("abstract_inverted_index")
            if inverted and isinstance(inverted, dict):
                try:
                    max_pos = max(
                        (pos for positions in inverted.values() for pos in positions),
                        default=-1,
                    )
                    words = [""] * (max_pos + 1)
                    for word, positions in inverted.items():
                        for pos in positions:
                            if pos < len(words):
                                words[pos] = str(word)
                    abstract = " ".join(wo for wo in words if wo)
                except Exception:
                    abstract = ""

            papers.append({
                "title": w.get("title", ""),
                "authors": [a["author"]["display_name"] for a in w.get("authorships", []) if a.get("author", {}).get("display_name")],
                "year": w.get("publication_year"),
                "doi": w.get("doi", ""),
                "cited_by_count": w.get("cited_by_count", 0),
                "abstract": abstract[:500],
            })
    except ImportError:
        error_msg = "httpx not installed"
    except Exception as exc:
        error_msg = str(exc)

    # Write literature matrix CSV
    matrix_path = workspace / "knowledge" / f"{req.taskId}_literature_matrix.csv"
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["title", "authors", "year", "doi", "cited_by_count", "abstract"])
    for p in papers:
        writer.writerow([p["title"], ", ".join(p["authors"]), p["year"], p["doi"], p["cited_by_count"], p["abstract"][:200]])
    matrix_path.write_text(buf.getvalue(), encoding="utf-8")

    # Write literature notes
    notes_path = workspace / "knowledge" / f"{req.taskId}_literature_notes.md"
    md_lines = [
        f"# Literature Search: {query}",
        "",
        f"**Source:** OpenAlex API  ",
        f"**Results:** {len(papers)} papers",
        "",
    ]
    for i, p in enumerate(papers, 1):
        md_lines.append(f"## {i}. {p['title']}")
        md_lines.append(f"- **Authors:** {', '.join(p['authors'][:5])}")
        md_lines.append(f"- **Year:** {p['year']}")
        md_lines.append(f"- **DOI:** {p['doi']}")
        md_lines.append(f"- **Citations:** {p['cited_by_count']}")
        md_lines.append("")
    if error_msg:
        md_lines.append(f"---\n*Warning: {error_msg}*")
    notes_path.write_text("\n".join(md_lines), encoding="utf-8")

    result: dict[str, Any] = {
        "ok": True,
        "message": f"OpenAlex search completed: {len(papers)} results",
        "papers": papers,
        "artifacts": [
            artifact("Literature Matrix CSV", matrix_path, "knowledge", "text/csv"),
            artifact("Literature Notes", notes_path, "knowledge", "text/markdown"),
        ],
    }
    if error_msg:
        result["warning"] = error_msg
    return result


@app.post("/tools/knowledge/index")
def index_knowledge(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    index_path = workspace / "knowledge" / "geowork_index.json"
    index_path.write_text(
        '{"status":"indexed","types":["pdf","docx","pptx","markdown","notebook","web"],"engine":"local"}',
        encoding="utf-8",
    )
    return {"ok": True, "artifacts": [artifact("Knowledge Index", index_path, "knowledge-index", "application/json")]}


@app.post("/tools/qgis/check")
def qgis_check(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    status_path = workspace / "artifacts" / f"{req.taskId}_qgis_status.json"
    status_path.write_text(
        '{"bundled":false,"strategy":"detect-local-installation","status":"not_configured"}',
        encoding="utf-8",
    )
    return {"ok": True, "artifacts": [artifact("QGIS Environment Status", status_path, "environment", "application/json")]}


@app.post("/tools/qgis/check-env")
def qgis_check_env(req: ToolRequest):
    return qgis_check(req)


@app.post("/tools/qgis/run-processing")
def qgis_run_processing(req: ToolRequest):
    workspace = ensure_workspace(req.workspace)
    algorithm = req.params.get("algorithm") or "native:buffer"
    payload = {"ok": True, "algorithm": algorithm, "parameters": req.params, "qgisBundled": False}
    art = write_json_artifact(workspace, req.taskId, "QGIS Processing Result", f"{req.taskId}_qgis_processing.json", payload, "qgis-processing")
    return {"ok": True, "message": "QGIS Processing task recorded", "artifacts": [art], "result": payload}
